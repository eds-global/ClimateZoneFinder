import streamlit as st
import base64
import pandas as pd
import plotly.express as px
import io
import re

st.set_page_config(
    page_title="Climate Analytics Dashboard",
    # page_icon="🌍",
    layout="wide"
)

def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode()

logo_base64 = get_base64_image("images/EDSlogo.jpg")

# STEP 1: Create header using st.columns for proper layout
col_logo, col_title, col_home = st.columns([1, 4, 1])

with col_logo:
    st.markdown(
        f'<img src="data:image/png;base64,{logo_base64}" style="height: 80px; margin-top: 45px;">',
        unsafe_allow_html=True
    )

with col_title:
    st.markdown(
        '<h2 style="text-align: center; color: #a85c42; margin-top: 45px;">Climate Analytics Dashboard</h2>',
        unsafe_allow_html=True
    )

with col_home:
    pass

def get_base64_image(image_path):
    try:
        with open(image_path, "rb") as img_file:
            return base64.b64encode(img_file.read()).decode()
    except:
        try:
            with open(f"../{image_path}", "rb") as img_file:
                return base64.b64encode(img_file.read()).decode()
        except:
            return ""

st.markdown("""
    <style>
    
    /* Style the header columns container */
    div[data-testid="stHorizontalBlock"]:first-of-type {
        border-bottom: 1px solid #e6e6e6;
        padding-bottom: 20px;
        margin-bottom: 0px;
        background-color: white;
    }
    

    
    /* Logo hover effect */
    div[data-testid="stHorizontalBlock"]:first-of-type img:hover {
        transform: scale(1.05);
        opacity: 0.85;
        transition: all 0.3s ease;
    }
    

    </style>
""", unsafe_allow_html=True)

# === HEADER ===
# st.markdown("""
#     <style>
#     .header-container {
#         background: linear-gradient(135deg, #1a3a52 0%, #2c5aa0 100%);
#         padding: 0px;
#         border-radius: 0;
#         margin-top: 50px;
#         box-shadow: 0 4px 12px rgba(0,0,0,0.2);
#         border-bottom: 4px solid #ffffff;

#            /* 👈 pushes it below Streamlit toolbar */
#         left: 0;
#         right: 0;
#         z-index: 999;
#         width: 100%;
#         box-sizing: border-box;
#     }
#     style>
#     /* Hide top toolbar */
#     header[data-testid="stHeader"] {
#         display: none;
#     }

#     /* Hide hamburger menu */
#     #MainMenu {
#         visibility: hidden;
#     }

#     /* Hide footer */
#     footer {
#         visibility: hidden;
#     }

#     /* Remove top padding since header is gone */
#     .block-container {
#         padding-top: 1rem;
#     }

#     /* Optional: Remove deploy button */
#     div[data-testid="stToolbar"] {
#         display: none;
#     }
#     /* Adjust body spacing to avoid overlap */
#     .main > div {
#         padding-top: 180px;   /* Increase if needed */
#     }

#     .header-content {
#         display: flex;
#         align-items: center;
#         gap: 20px;
#     }

#     .header-icon {
#         font-size: 48px;
#         display: inline-block;
#     }

#     .header-title {
#         color: #ffffff;
#         font-size: 32px;
#         font-weight: 800;
#         margin: 0;
#         letter-spacing: 0.5px;
#     }
#     </style>

#     <div class="header-container">
#         <div class="header-content">
#             <div class="header-icon">🌍</div>
#             <div class="header-title">Climate Analytics Dashboard</div>
#         </div>
#     </div>
# """, unsafe_allow_html=True)

# === INITIALIZE SESSION STATE FOR TABS ===
if "active_tab" not in st.session_state:
    st.session_state.active_tab = "Annual Trend"

st.markdown("""
    <style>
    /* Main layout adjustments */
    .block-container {
        padding-top: 0rem !important;
    }
    
    section[data-testid="stSidebar"] {
        display: none !important;
    }
    
    button[kind="header"] {
        display: none !important;
    }
    
    .main .block-container {
        max-width: 100% !important;
        padding-left: 1.5rem !important;
        padding-right: 1.5rem !important;
    }
    
    
    
    .control-section-header {
        font-size: 15px;
        font-weight: 700;
        color: #2c3e50;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        margin-bottom: 12px;
        margin-top: 16px;
        display: flex;
        align-items: center;
        gap: 8px;
        width: 200px;
    }
    
    .control-section-header:first-child {
        margin-top: 0;
    }
    
    /* Upload zone styling */
    .upload-zone {
        border: 2px dashed #cbd5e0;
        border-radius: 6px;
        padding: 12px;
        text-align: center;
        background-color: #f7fafc;
        margin-top: 8px;
    }
    
    .upload-zone.success {
        border-color: #48bb78;
        background-color: #f0fff4;
    }
    
    /* File uploader styling */
    [data-testid="fileUploadDropzone"] {
        border-radius: 6px !important;
        border-color: #cbd5e0 !important;
    }
    
    /* Success message styling */
    .stAlert[data-baseweb="notification"] {
        background-color: #f0fff4 !important;
        border-left: 4px solid #48bb78 !important;
        border-radius: 4px !important;
    }
    
    /* Slider styling */
    .stSlider {
        margin-top: 12px;
    }
    
    /* Date input styling */
    .stDateInput {
        margin-top: 12px;
    }
    
    /* Selectbox styling */
    .stSelectbox {
        margin-top: 8px;
    }
    
    /* Section styling */
    .section-title {
        font-size: 18px;
        font-weight: 700;
        color: #2c3e50;
        margin-bottom: 16px;
        border-bottom: 2px solid #3498db;
        padding-bottom: 8px;
        display: inline-block;
    }
    
    /* KPI Cards */
    .kpi-card {
        background: white;
        padding: 16px;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
        text-align: center;
    }
    
    .kpi-label {
        font-size: 11px;
        font-weight: 700;
        color: #718096;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        margin-bottom: 8px;
        opacity: 0.9;
    }
    
    .kpi-value {
        font-size: 26px;
        font-weight: 700;
        color: #2c3e50;
        margin: 8px 0;
    }
    
    .kpi-meta {
        font-size: 11px;
        color: #718096;
        opacity: 0.85;
    }
    </style>
""", unsafe_allow_html=True)

def generate_pptx_report(df: pd.DataFrame, start_date, end_date, start_hour: int, end_hour: int, selected_parameter: str):
    """Generate a PowerPoint report with Dry Bulb and Humidity analysis."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import io
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    import tempfile
    import os
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define a color scheme
    DARK_BLUE = RGBColor(44, 90, 160)  # #2c5aa0
    LIGHT_BLUE = RGBColor(26, 58, 82)  # #1a3a52
    ACCENT_RED = RGBColor(211, 47, 47)  # #d32f2f
    TEXT_COLOR = RGBColor(44, 62, 80)  # #2c3e50
    
    def add_title_slide(prs, title, subtitle):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        
        return slide
    
    def add_content_slide(prs, title, content_func):
        """Add a content slide with title."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = LIGHT_BLUE
        title_shape.line.color.rgb = DARK_BLUE
        
        # Add title text
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Call content function to add content
        content_func(slide)
        
        return slide
    
    # Filter data for period
    filtered_df = df[
        (df["datetime"].dt.date >= start_date) &
        (df["datetime"].dt.date <= end_date) &
        (df["hour"].between(start_hour, end_hour))
    ]
    
    # === TITLE SLIDE ===
    add_title_slide(
        prs, 
        "Climate Analysis Report",
        f"Period: {start_date.strftime('%b %d')} - {end_date.strftime('%b %d')}\nHours: {start_hour:02d}:00 - {end_hour:02d}:00"
    )
    
    # === ANNUAL TREND CHART SLIDE (DRY BULB) ===
    def add_temp_chart_content(slide):
        try:
            # Compute daily stats
            daily_stats = df.groupby("doy").agg({
                "dry_bulb_temperature": ["min", "max", "mean"],
            }).reset_index()
            
            daily_stats.columns = ["doy", "temp_min", "temp_max", "temp_avg"]
            
            # Calculate ASHRAE comfort bands
            daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
            comfort_line = daily_avg.rolling(window=7, center=True).mean()
            comfort_80_lower = comfort_line - 3.5
            comfort_80_upper = comfort_line + 3.5
            comfort_90_lower = comfort_line - 2.5
            comfort_90_upper = comfort_line + 2.5
            
            # Create matplotlib figure
            fig, ax = plt.subplots(figsize=(12, 5.5), dpi=120)
            
            # Calculate date range for filtering
            start_month_num = start_date.month
            end_month_num = end_date.month
            start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
            if end_month_num == 12:
                end_doy = 366
            else:
                end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
            
            # Plot ASHRAE comfort bands
            ax.fill_between(daily_stats["doy"], comfort_80_lower, comfort_80_upper, 
                           alpha=0.2, color='gray', label='ASHRAE 80% Comfort')
            ax.fill_between(daily_stats["doy"], comfort_90_lower, comfort_90_upper, 
                           alpha=0.3, color='gray', label='ASHRAE 90% Comfort')
            
            # Plot temperature range
            ax.fill_between(daily_stats["doy"], daily_stats["temp_min"], daily_stats["temp_max"],
                           alpha=0.35, color='#FFB3B3', label='Daily Temp Range')
            
            # Plot average line
            ax.plot(daily_stats["doy"], daily_stats["temp_avg"], 
                   color='#d32f2f', linewidth=2.5, label='Daily Average', zorder=3)
            
            # Highlight selected period
            ax.axvspan(start_doy, end_doy, alpha=0.08, color='#2c5aa0')
            
            ax.set_xlabel('Day of Year', fontsize=11, fontweight='bold')
            ax.set_ylabel('Temperature (°C)', fontsize=11, fontweight='bold')
            ax.set_title('Annual Dry Bulb Temperature Trend', fontsize=13, fontweight='bold', pad=15)
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.12), ncol=4, frameon=True, fontsize=10)
            ax.grid(True, alpha=0.3, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            
            plt.tight_layout()
            
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                fig.savefig(tmp.name, dpi=120, bbox_inches='tight', facecolor='white')
                tmp_path = tmp.name
            
            plt.close(fig)
            
            # Add image to slide
            slide.shapes.add_picture(tmp_path, Inches(0.3), Inches(1.0), width=Inches(9.4), height=Inches(5.8))
            
            # Clean up
            os.unlink(tmp_path)
            
        except Exception as e:
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = f"Chart visualization error: {str(e)}"
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    add_content_slide(prs, "Annual Trend - Dry Bulb Temperature", add_temp_chart_content)
    
    # === ANNUAL TREND CHART SLIDE (HUMIDITY) ===
    def add_humidity_chart_content(slide):
        try:
            # Compute daily stats for humidity
            daily_stats = df.groupby("doy").agg({
                "relative_humidity": ["min", "max", "mean"],
            }).reset_index()
            
            daily_stats.columns = ["doy", "rh_min", "rh_max", "rh_avg"]
            
            # Create matplotlib figure
            fig, ax = plt.subplots(figsize=(12, 5.5), dpi=120)
            
            # Calculate date range for filtering
            start_month_num = start_date.month
            end_month_num = end_date.month
            start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
            if end_month_num == 12:
                end_doy = 366
            else:
                end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
            
            # Plot risk zones
            ax.axhspan(75, 100, alpha=0.15, color='#FF6B6B', label='Condensation Risk (>75%)')
            ax.axhspan(60, 75, alpha=0.15, color='#FFA500', label='High RH (60-75%)')
            ax.axhspan(30, 60, alpha=0.15, color='#4ECDC4', label='Comfortable (30-60%)')
            ax.axhspan(0, 30, alpha=0.15, color='#FFD93D', label='Low RH (<30%)')
            
            # Plot humidity range
            ax.fill_between(daily_stats["doy"], daily_stats["rh_min"], daily_stats["rh_max"],
                           alpha=0.35, color='#0099ff', label='Daily RH Range')
            
            # Plot average line
            ax.plot(daily_stats["doy"], daily_stats["rh_avg"], 
                   color='#0066cc', linewidth=2.5, label='Daily Average RH', zorder=3)
            
            # Highlight selected period
            ax.axvspan(start_doy, end_doy, alpha=0.08, color='#2c5aa0')
            
            ax.set_xlabel('Day of Year', fontsize=11, fontweight='bold')
            ax.set_ylabel('Relative Humidity (%)', fontsize=11, fontweight='bold')
            ax.set_title('Annual Relative Humidity Trend', fontsize=13, fontweight='bold', pad=15)
            ax.set_ylim(0, 100)
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.12), ncol=4, frameon=True, fontsize=10)
            ax.grid(True, alpha=0.3, linestyle='--')
            ax.set_facecolor('#fafafa')
            fig.patch.set_facecolor('white')
            
            plt.tight_layout()
            
            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                fig.savefig(tmp.name, dpi=120, bbox_inches='tight', facecolor='white')
                tmp_path = tmp.name
            
            plt.close(fig)
            
            # Add image to slide
            slide.shapes.add_picture(tmp_path, Inches(0.3), Inches(1.0), width=Inches(9.4), height=Inches(5.8))
            
            # Clean up
            os.unlink(tmp_path)
            
        except Exception as e:
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = f"Chart visualization error: {str(e)}"
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    add_content_slide(prs, "Annual Trend - Relative Humidity", add_humidity_chart_content)
    
    # === DRY BULB TEMPERATURE ANALYSIS ===
    def add_temp_content(slide):
        if not filtered_df.empty:
            min_temp = filtered_df["dry_bulb_temperature"].min()
            max_temp = filtered_df["dry_bulb_temperature"].max()
            avg_temp = filtered_df["dry_bulb_temperature"].mean()
            
            # Calculate HDD18 and CDD24 for full year
            hdd18 = (18 - df["dry_bulb_temperature"]).clip(lower=0).sum()
            cdd24 = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
            
            hdd18_filtered = (18 - filtered_df["dry_bulb_temperature"]).clip(lower=0).sum()
            cdd24_filtered = (filtered_df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
            
            # Create metrics with better layout - 2 rows of 4 cards
            metrics = [
                ("Min Temp", f"{min_temp:.1f}°C", "#FF9800", RGBColor(255, 152, 0)),
                ("Max Temp", f"{max_temp:.1f}°C", "#F44336", RGBColor(244, 67, 54)),
                ("Avg Temp", f"{avg_temp:.1f}°C", "#2196F3", RGBColor(33, 150, 243)),
                ("Diurnal Range", f"{max_temp - min_temp:.1f}°C", "#9C27B0", RGBColor(156, 39, 176)),
                ("HDD18 (Annual)", f"{hdd18:.0f}", "#1976D2", RGBColor(25, 118, 210)),
                ("CDD24 (Annual)", f"{cdd24:.0f}", "#D32F2F", RGBColor(211, 47, 47)),
                ("HDD18 (Period)", f"{hdd18_filtered:.0f}", "#0097A7", RGBColor(0, 150, 167)),
                ("CDD24 (Period)", f"{cdd24_filtered:.0f}", "#C62828", RGBColor(198, 40, 40)),
            ]
            
            for idx, (label, value, hex_color, rgb_color) in enumerate(metrics):
                col = idx % 4
                row = idx // 4
                
                left = Inches(0.4 + col * 2.35)
                top = Inches(1.1 + row * 1.6)
                box_width = Inches(2.15)
                box_height = Inches(1.4)
                
                # Add rounded rectangle box with gradient-like effect using border
                box = slide.shapes.add_shape(1, left, top, box_width, box_height)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                box.line.color.rgb = rgb_color
                box.line.width = Pt(3)
                
                # Add colored top border for visual appeal
                top_bar = slide.shapes.add_shape(1, left, top, box_width, Inches(0.25))
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = rgb_color
                top_bar.line.color.rgb = rgb_color
                
                # Add label
                label_frame = box.text_frame
                label_frame.clear()
                label_frame.word_wrap = True
                label_frame.margin_top = Inches(0.08)
                label_frame.margin_left = Inches(0.1)
                label_frame.margin_right = Inches(0.1)
                
                # Label
                p = label_frame.paragraphs[0]
                p.text = label
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = rgb_color
                
                # Value
                p2 = label_frame.add_paragraph()
                p2.text = value
                p2.font.size = Pt(16)
                p2.font.bold = True
                p2.font.color.rgb = TEXT_COLOR
                p2.space_before = Pt(4)
    
    add_content_slide(prs, "Dry Bulb Temperature Analysis", add_temp_content)
    
    # === HUMIDITY ANALYSIS ===
    def add_humidity_content(slide):
        if not filtered_df.empty:
            # Full year metrics
            high_humidity_annual = len(df[df["relative_humidity"] > 60])
            condensation_risk_annual = len(df[df["relative_humidity"] > 75])
            low_humidity_annual = len(df[df["relative_humidity"] < 30])
            avg_rh_annual = df["relative_humidity"].mean()
            
            # Period metrics
            min_rh = filtered_df["relative_humidity"].min()
            max_rh = filtered_df["relative_humidity"].max()
            avg_rh = filtered_df["relative_humidity"].mean()
            high_humidity_filtered = len(filtered_df[filtered_df["relative_humidity"] > 60])
            condensation_risk_filtered = len(filtered_df[filtered_df["relative_humidity"] > 75])
            
            # Create metrics with better layout
            metrics = [
                ("Min RH", f"{min_rh:.0f}%", "#4CAF50", RGBColor(76, 175, 80)),
                ("Max RH", f"{max_rh:.0f}%", "#FF5722", RGBColor(255, 87, 34)),
                ("Avg RH", f"{avg_rh:.0f}%", "#00BCD4", RGBColor(0, 188, 212)),
                ("Low RH (<30%)", f"{low_humidity_annual:.0f} hrs", "#FFB74D", RGBColor(255, 183, 77)),
                ("High RH (>60%)", f"{high_humidity_annual:.0f} hrs", "#EF5350", RGBColor(239, 83, 80)),
                ("Condensation (>75%)", f"{condensation_risk_annual:.0f} hrs", "#E53935", RGBColor(229, 57, 53)),
                ("High RH (Period)", f"{high_humidity_filtered:.0f} hrs", "#AB47BC", RGBColor(171, 71, 188)),
                ("Condensation (Period)", f"{condensation_risk_filtered:.0f} hrs", "#8E24AA", RGBColor(142, 36, 170)),
            ]
            
            for idx, (label, value, hex_color, rgb_color) in enumerate(metrics):
                col = idx % 4
                row = idx // 4
                
                left = Inches(0.4 + col * 2.35)
                top = Inches(1.1 + row * 1.6)
                box_width = Inches(2.15)
                box_height = Inches(1.4)
                
                # Add rounded rectangle box
                box = slide.shapes.add_shape(1, left, top, box_width, box_height)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                box.line.color.rgb = rgb_color
                box.line.width = Pt(3)
                
                # Add colored top border
                top_bar = slide.shapes.add_shape(1, left, top, box_width, Inches(0.25))
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = rgb_color
                top_bar.line.color.rgb = rgb_color
                
                # Add label
                label_frame = box.text_frame
                label_frame.clear()
                label_frame.word_wrap = True
                label_frame.margin_top = Inches(0.08)
                label_frame.margin_left = Inches(0.1)
                label_frame.margin_right = Inches(0.1)
                
                # Label
                p = label_frame.paragraphs[0]
                p.text = label
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = rgb_color
                
                # Value
                p2 = label_frame.add_paragraph()
                p2.text = value
                p2.font.size = Pt(16)
                p2.font.bold = True
                p2.font.color.rgb = TEXT_COLOR
                p2.space_before = Pt(4)
    
    add_content_slide(prs, "Humidity Analysis", add_humidity_content)
    
    # === SUMMARY SLIDE ===
    def add_summary_content(slide):
        summary_text = f"""
Period Analyzed: {start_date.strftime('%B %d')} - {end_date.strftime('%B %d')}
Operating Hours: {start_hour:02d}:00 - {end_hour:02d}:00

Dry Bulb Temperature:
• Range: {filtered_df['dry_bulb_temperature'].min():.2f}°C to {filtered_df['dry_bulb_temperature'].max():.2f}°C
• Average: {filtered_df['dry_bulb_temperature'].mean():.2f}°C
• Hours above 28°C: {len(filtered_df[filtered_df['dry_bulb_temperature'] > 28])}
• Hours below 12°C: {len(filtered_df[filtered_df['dry_bulb_temperature'] < 12])}

Relative Humidity:
• Range: {filtered_df['relative_humidity'].min():.1f}% to {filtered_df['relative_humidity'].max():.1f}%
• Average: {filtered_df['relative_humidity'].mean():.1f}%
• Hours with High RH (>60%): {len(filtered_df[filtered_df['relative_humidity'] > 60])}
• Hours with Condensation Risk (>75%): {len(filtered_df[filtered_df['relative_humidity'] > 75])}
        """
        
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = summary_text
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT_COLOR
            if paragraph.text.startswith("•"):
                paragraph.level = 1
                paragraph.font.size = Pt(11)
            elif ":" in paragraph.text and not paragraph.text.startswith("•"):
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
    
    add_content_slide(prs, "Summary", add_summary_content)
    
    # Save to bytes
    report_bytes = io.BytesIO()
    prs.save(report_bytes)
    report_bytes.seek(0)
    
    return report_bytes

def parse_epw(epw_text: str) -> pd.DataFrame:
    """Parse EPW formatted text and return a DataFrame with datetime, dry_bulb_temperature, relative_humidity."""
    # split into lines and find first data row (starts with a year integer)
    lines = [ln.strip() for ln in epw_text.splitlines() if ln.strip() != ""]
    data_start = None
    for i, ln in enumerate(lines):
        # consider a line a data line if first token is a 4-digit year
        toks = ln.split(",")
        if len(toks) > 1 and re.fullmatch(r"\d{4}", toks[0].strip()):
            data_start = i
            break
    
    if data_start is None:
        raise ValueError("Could not locate EPW data rows")
    
    data_str = "\n".join(lines[data_start:])
    df_raw = pd.read_csv(io.StringIO(data_str), header=None)
    
    # EPW standard columns (0-based index):
    # 0=year,1=month,2=day,3=hour,4=minute,5=data source,6=dry bulb (C),7=dew point,8=relative humidity (%)
    col_map = {
        "year": 0,
        "month": 1,
        "day": 2,
        "hour": 3,
        "minute": 4,
        "dry_bulb_temperature": 6,
        "relative_humidity": 8,
    }
    
    # safety: ensure DataFrame has enough columns
    max_needed = max(col_map.values())
    if df_raw.shape[1] <= max_needed:
        raise ValueError("EPW data appears to have insufficient columns")
    
    df = pd.DataFrame()
    df["year"] = pd.to_numeric(df_raw.iloc[:, col_map["year"]], errors="coerce").astype("Int64")
    df["month"] = pd.to_numeric(df_raw.iloc[:, col_map["month"]], errors="coerce").astype("Int64")
    df["day"] = pd.to_numeric(df_raw.iloc[:, col_map["day"]], errors="coerce").astype("Int64")
    df["hour_raw"] = pd.to_numeric(df_raw.iloc[:, col_map["hour"]], errors="coerce").astype("Int64")
    df["minute"] = pd.to_numeric(df_raw.iloc[:, col_map["minute"]], errors="coerce").astype("Int64")
    
    # EPW hours are 1-24 representing the hour ending; convert to 0-23 by subtracting 1
    df["hour"] = (df["hour_raw"].fillna(1).astype(int) - 1) % 24
    
    df["dry_bulb_temperature"] = pd.to_numeric(df_raw.iloc[:, col_map["dry_bulb_temperature"]], errors="coerce")
    df["relative_humidity"] = pd.to_numeric(df_raw.iloc[:, col_map["relative_humidity"]], errors="coerce")
    
    # build datetime (note: this may produce NaT for invalid rows)
    df["datetime"] = pd.to_datetime(
        dict(year=df["year"], month=df["month"], day=df["day"], hour=df["hour"], minute=df["minute"]),
        errors="coerce",
    )
    
    df = df.dropna(subset=["datetime"]).reset_index(drop=True)
    return df[["datetime", "dry_bulb_temperature", "relative_humidity", "hour"]]

# === MAIN LAYOUT ===
col_left, col_right = st.columns([0.85, 2.15], gap="large")

with col_left:
    # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
    
    # Upload EPW File Section
    # st.markdown('<div class="control-section-header">📤 Upload EPW File</div>', unsafe_allow_html=True)
    # st.markdown('<div class="upload-zone">Limit 200MB per file · EPW</div>', unsafe_allow_html=True)
    st.write("##### 📤 Upload EPW File")
    uploaded = st.file_uploader("", type=["epw"], label_visibility="collapsed", width=300)
    
    # Parameter Selection
    # st.markdown('<div class="control-section-header">⚙️ Parameter</div>', unsafe_allow_html=True)
    st.write("##### ⚙️ Parameter")
    selected_parameter = st.selectbox(
        "Select parameter",
        ["Temperature", "Humidity", "Sun Path"],
        label_visibility="collapsed",
        key="parameter_selector",
        width=300
    )

if uploaded is None:
    with col_right:
        st.info("Please upload an .epw file to analyze.", width=300)
    st.stop()

try:
    raw = uploaded.getvalue().decode("utf-8", errors="replace")
    df = parse_epw(raw)
    
    # compute derived date fields used by the charts
    df["doy"] = df["datetime"].dt.dayofyear
    df["day"] = df["datetime"].dt.day
    df["month"] = df["datetime"].dt.month
    df["month_name"] = df["datetime"].dt.strftime("%b")
    
    with col_left:
        st.markdown("""
            <div style="
                background-color: #f0fff4;
                border-left: 4px solid #48bb78;
                padding: 12px;
                border-radius: 4px;
                margin: 8px 0;
                width: 300px;
            ">
                <div style="color: #22543d; font-weight: 600; font-size: 12px;">✅ EPW parsed successfully</div>
            </div>
        """, unsafe_allow_html=True)
        
        # Time range (hour of use) slider for diurnal/peak analysis
        st.markdown('<div class="control-section-header">⏰ Time Range (Hours)</div>', unsafe_allow_html=True)
        hour_range = st.slider(
            "Select hours (start - end)",
            min_value=0,
            max_value=23,
            value=(8, 18),
            step=1,
            key="hour_range",
            label_visibility="collapsed",
            width=300
        )
        
        # Date range selection
        st.markdown('<div class="control-section-header">📅 Date Range</div>', unsafe_allow_html=True)
        
        months_list = ["January", "February", "March", "April", "May", "June", 
                       "July", "August", "September", "October", "November", "December"]
        
        # Initialize session state for month selection
        if "start_month_idx" not in st.session_state:
            st.session_state.start_month_idx = 0
        if "end_month_idx" not in st.session_state:
            st.session_state.end_month_idx = 11
        
        # Create two columns for start and end month dropdowns
        month_col1, month_col2, col3 = st.columns([1,1,0.3], gap="small")
        
        with month_col1:
            start_month = st.selectbox(
                "From",
                options=range(len(months_list)),
                format_func=lambda x: months_list[x],
                key="start_month_select",
                label_visibility="collapsed",
                width=150
            )
            st.session_state.start_month_idx = start_month
        
        with month_col2:
            # End month should only show months from start_month onwards
            end_month_options = list(range(start_month, len(months_list)))
            # Ensure selected end_month is valid
            if st.session_state.end_month_idx < start_month:
                st.session_state.end_month_idx = start_month
            
            end_month = st.selectbox(
                "To",
                options=end_month_options,
                format_func=lambda x: months_list[x],
                key="end_month_select",
                index=min(st.session_state.end_month_idx - start_month, len(end_month_options) - 1),
                label_visibility="collapsed",
                width=150,
            )
            st.session_state.end_month_idx = end_month
        
        with col3:
            pass
        
        # Generate Report Button
        st.markdown('<div class="control-section-header">📊 Report(PowerPoint)</div>', unsafe_allow_html=True)
        if st.button("Generate Report", use_container_width=False, key="generate_report_btn", width=300):
            st.session_state.generate_report = True
        
except Exception as e:
    with col_left:
        st.error(f"❌ Failed to parse EPW: {e}")
    st.stop()

def calculate_ashrae_comfort(df: pd.DataFrame) -> tuple:
    """
    Calculate ASHRAE adaptive comfort bands.
    Returns (comfort_80_lower, comfort_80_upper, comfort_90_lower, comfort_90_upper) as daily rolling averages.
    """
    # Simple ASHRAE adaptive comfort model (simplified)
    # For illustration: ranges relative to outdoor mean monthly temperature
    daily_avg = df.groupby("doy")["dry_bulb_temperature"].mean()
    
    # 80% acceptability: ±3.5°C from comfort line
    # 90% acceptability: ±2.5°C from comfort line
    comfort_line = daily_avg.rolling(window=7, center=True).mean()
    
    comfort_80_lower = comfort_line - 3.5
    comfort_80_upper = comfort_line + 3.5
    comfort_90_lower = comfort_line - 2.5
    comfort_90_upper = comfort_line + 2.5
    
    return comfort_80_lower, comfort_80_upper, comfort_90_lower, comfort_90_upper

# Compute ASHRAE comfort bands
ashrae_80_lower, ashrae_80_upper, ashrae_90_lower, ashrae_90_upper = calculate_ashrae_comfort(df)

with col_right:
    # === INTERACTIVE TABS AT TOP ===
    st.markdown("""
        <style>
        .tab-container {
            display: flex;
            gap: 0;
            background-color: #f8f9fa;
            padding: 0;
            margin: -1rem -1rem 1.5rem -1rem;
            border-bottom: 2px solid #e9ecef;
        }
        .tab-button {
            padding: 12px 24px;
            cursor: pointer;
            font-size: 14px;
            font-weight: 600;
            color: #495057;
            background-color: #f8f9fa;
            border: none;
            border-bottom: 3px solid transparent;
            transition: all 0.3s ease;
        }
        .tab-button:hover {
            background-color: #e9ecef;
        }
        .tab-button.active {
            color: #2c3e50;
            border-bottom-color: #3498db;
            background-color: white;
        }
        </style>
    """, unsafe_allow_html=True)
    
    # Create tab buttons
    tabs_col1, tabs_col2, tabs_col3, tabs_col4, tabs_col5 = st.columns(5, gap="small")
    
    with tabs_col1:
        if st.button("Annual Trend", key="tab_annual", use_container_width=True):
            st.session_state.active_tab = "Annual Trend"
    
    with tabs_col2:
        if st.button("Monthly Trend", key="tab_monthly", use_container_width=True):
            st.session_state.active_tab = "Monthly Trend"
    
    with tabs_col3:
        if st.button("Diurnal Profile", key="tab_diurnal", use_container_width=True):
            st.session_state.active_tab = "Diurnal Profile"
    
    with tabs_col4:
        if st.button("Comfort Analysis", key="tab_comfort", use_container_width=True):
            st.session_state.active_tab = "Comfort Analysis"
    
    with tabs_col5:
        if st.button("Energy Metrics", key="tab_energy", use_container_width=True):
            st.session_state.active_tab = "Energy Metrics"
    
    # Add visual styling to show active tab
    st.markdown(f"""
        <script>
        var active_tab = '{st.session_state.active_tab}';
        </script>
    """, unsafe_allow_html=True)
    # Get the actual year from the data
    year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
    
    # Create start and end dates based on selected months
    start_month_num = st.session_state.start_month_idx + 1
    end_month_num = st.session_state.end_month_idx + 1
    
    start_date = pd.to_datetime(f"{year}-{start_month_num}-01").date()
    
    # For end_date, get the last day of the end month
    if end_month_num == 12:
        end_date = pd.to_datetime(f"{year}-12-31").date()
    else:
        end_date = (pd.to_datetime(f"{year}-{end_month_num+1}-01") - pd.Timedelta(days=1)).date()
    
    start_hour, end_hour = st.session_state.get("hour_range", (8, 18))
    
    # Compute daily min/max and average
    daily_stats = df.groupby("doy").agg({
        "dry_bulb_temperature": ["min", "max", "mean"],
        "relative_humidity": ["min", "max", "mean"],
    }).reset_index()
    
    daily_stats.columns = ["doy", "temp_min", "temp_max", "temp_avg", "rh_min", "rh_max", "rh_avg"]
    daily_stats["datetime_start"] = pd.to_datetime(daily_stats["doy"], format="%j", errors="coerce")
    
    # Map doy back to actual year-dates for chart display
    year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
    daily_stats["datetime"] = pd.to_datetime(
        daily_stats["doy"].astype(str) + f"-{year}", format="%j-%Y", errors="coerce"
    )
    
    # Add a day-month only column for display (without year)
    daily_stats["datetime_display"] = daily_stats["datetime"].dt.strftime("%b %d")
    
    # Map comfort bands to doy
    comfort_df = pd.DataFrame({
        "doy": ashrae_80_lower.index,
        "comfort_80_lower": ashrae_80_lower.values,
        "comfort_80_upper": ashrae_80_upper.values,
        "comfort_90_lower": ashrae_90_lower.values,
        "comfort_90_upper": ashrae_90_upper.values,
    })
    
    # Merge comfort data with daily stats
    daily_stats = daily_stats.merge(comfort_df, on="doy", how="left")

    # === TAB 1: ANNUAL TREND ===
    if st.session_state.active_tab == "Annual Trend":
        if selected_parameter == "Temperature":
            import plotly.graph_objects as go
            
            # Calculate day of year for selected month range for greying
            start_month_num = st.session_state.start_month_idx + 1
            end_month_num = st.session_state.end_month_idx + 1
            start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
            if end_month_num == 12:
                end_doy = 366  # Use 366 to catch Dec 31
            else:
                end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
            
            fig_yearly = go.Figure()
            
            # GREYED OUT: Before selected range
            if start_doy > 1:
                before_data = daily_stats[daily_stats["doy"] < start_doy]
                fig_yearly.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_yearly.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Period",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                fig_yearly.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_avg"],
                    mode="lines",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected range
            active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
            # Add ASHRAE 80% band for ACTIVE range
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_80_upper"],
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_80_lower"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="ASHRAE adaptive comfort (80%)",
                fillcolor="rgba(128, 128, 128, 0.2)",
                hoverinfo="skip",
            ))
            
            # Add ASHRAE 90% band
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_90_upper"],
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_90_lower"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="ASHRAE adaptive comfort (90%)",
                fillcolor="rgba(128, 128, 128, 0.4)",
                hoverinfo="skip",
            ))
            
            # Add temperature range (min/max) for ACTIVE range
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_max"],
                fill=None,
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                name="Dry bulb temperature Range",
                fillcolor="rgba(255, 173, 173, 0.4)",
                customdata=active_range["temp_max"],
                hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
            ))
            
            # Add average line for ACTIVE range
            fig_yearly.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_avg"],
                mode="lines",
                name="Average Dry bulb temperature",
                line=dict(color="#d32f2f", width=2),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
            ))
            
            # GREYED OUT: After selected range
            if end_doy < 365:
                after_data = daily_stats[daily_stats["doy"] > end_doy]
                if not after_data.empty:
                    fig_yearly.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_yearly.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_yearly.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
            
            
            fig_yearly.update_layout(
                title="Annual Dry Bulb Temperature Trend",
                xaxis_title=None,
                yaxis_title="Temperature (°C)",
                hovermode="x unified",
                showlegend=True,
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=1.02,
                    xanchor="right",
                    x=1
                ),
                xaxis_rangeslider_visible=False,
                height=450,
                template="plotly_white",
                margin=dict(b=80),
            )

            st.plotly_chart(fig_yearly, use_container_width=True)
            
            # Get filtered data for period
            filtered_df = df[
                (df["datetime"].dt.date >= start_date) &
                (df["datetime"].dt.date <= end_date) &
                (df["hour"].between(start_hour, end_hour))
            ]
            
            if not filtered_df.empty:
                # Get min/max values and their corresponding rows with datetime/hour info
                min_idx = filtered_df["dry_bulb_temperature"].idxmin()
                max_idx = filtered_df["dry_bulb_temperature"].idxmax()
                
                min_row = filtered_df.loc[min_idx]
                max_row = filtered_df.loc[max_idx]
                
                temp_min = min_row["dry_bulb_temperature"]
                temp_max = max_row["dry_bulb_temperature"]
                temp_avg = filtered_df["dry_bulb_temperature"].mean()
                diurnal_range = temp_max - temp_min
                
                # Extract date and hour information for min/max
                min_date_str = min_row["datetime"].strftime("%b %d")
                min_hour = int(min_row["hour"])
                
                max_date_str = max_row["datetime"].strftime("%b %d")
                max_hour = int(max_row["hour"])
                
                # Calculate additional metrics using full year data
                # HDD18 (Heating Degree Days at 18°C base)
                hdd18 = (18 - df["dry_bulb_temperature"]).clip(lower=0).sum()
                
                # CDD24 (Cooling Degree Days at 24°C base)
                cdd24 = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                
                # Comfort metrics (using full year data for comfort bands)
                def get_comfort_band_range(temps):
                    """Get comfort band as ±3.5°C from mean"""
                    mean_temp = temps.mean()
                    return mean_temp - 3.5, mean_temp + 3.5
                
                comfort_lower, comfort_upper = get_comfort_band_range(df["dry_bulb_temperature"])
                comfort_hours = len(df[(df["dry_bulb_temperature"] >= comfort_lower) & (df["dry_bulb_temperature"] <= comfort_upper)])
                comfort_80_percent = (comfort_hours / len(df)) * 100
                
                # 1% Cooling (99th percentile)
                cooling_1pct = df["dry_bulb_temperature"].quantile(0.99)
                
                # Overheat hours (hours above 28°C threshold)
                overheat_hrs = len(df[df["dry_bulb_temperature"] > 28])
                
                # Cold hours (hours below 12°C threshold)
                cold_hrs = len(df[df["dry_bulb_temperature"] < 12])
                
                # Temperature metrics cards - Row 1
                kpi_col1, kpi_col2, kpi_col3, kpi_col4, kpi_col5 = st.columns(5)
                
                with kpi_col1:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #f59e0b;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Min Temp</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_min:.2f} °C</div>
                            <div style="font-size: 11px; color: #718096;">{min_date_str} · {min_hour:02d}:00</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col2:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #ef4444;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Max Temp</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_max:.2f} °C</div>
                            <div style="font-size: 11px; color: #718096;">{max_date_str} · {max_hour:02d}:00</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col3:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #8b5cf6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">Avg Temp</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{temp_avg:.2f} °C</div>
                            <div style="font-size: 11px; color: #718096;">All year average</div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col4:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #3b82f6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Diurnal Range</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{diurnal_range:.2f} °C</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col5:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #06b6d4;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">1% Cooling</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cooling_1pct:.2f} °C</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                # Row 2 - Additional metrics
                kpi_col6, kpi_col7, kpi_col8, kpi_col9, kpi_col10 = st.columns(5)
                
                with kpi_col6:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #dc2626;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #dc2626; text-transform: uppercase; letter-spacing: 0.5px;">HDD18</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{hdd18:.0f}</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col7:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #0891b2;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #0891b2; text-transform: uppercase; letter-spacing: 0.5px;">CDD24</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cdd24:.0f}</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col8:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #06b6d4;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">Comfort 80%</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{comfort_80_percent:.0f} %</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col9:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #8b5cf6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">Overheat Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{overheat_hrs}</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
                
                with kpi_col10:
                    st.markdown(f"""
                        <div style="
                            background: white;
                            padding: 16px;
                            border-radius: 8px;
                            border-left: 4px solid #3b82f6;
                            box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                            text-align: center;
                        ">
                            <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Cold Hrs</div>
                            <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{cold_hrs}</div>
                            <div style="font-size: 11px; color: #718096;"></div>
                        </div>
                    """, unsafe_allow_html=True)
            else:
                st.info(f"No data available between {start_hour:02d}:00 and {end_hour:02d}:00 in the selected date range.")
        
        elif selected_parameter == "Humidity":
            import plotly.graph_objects as go
            
            # Define humidity comfort bands (typical comfort range 30-65%)
            humidity_comfort_lower = 30
            humidity_comfort_upper = 65
            
            fig_yearly = go.Figure()
            
            # Add humidity comfort band
            fig_yearly.add_trace(go.Scatter(
                x=daily_stats["datetime_display"],
                y=[humidity_comfort_upper] * len(daily_stats),
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_yearly.add_trace(go.Scatter(
                x=daily_stats["datetime_display"],
                y=[humidity_comfort_lower] * len(daily_stats),
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="Humidity comfort band",
                fillcolor="rgba(128, 128, 128, 0.2)",
                hoverinfo="skip",
            ))
            
            # Add relative humidity range (min/max)
            fig_yearly.add_trace(go.Scatter(
                x=daily_stats["datetime_display"],
                y=daily_stats["rh_max"],
                fill=None,
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_yearly.add_trace(go.Scatter(
                x=daily_stats["datetime_display"],
                y=daily_stats["rh_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                name="Relative humidity Range",
                fillcolor="rgba(0, 150, 255, 0.3)",
                hovertemplate="<b>%{x}</b><br>Min: %{y:.1f}%<extra></extra>",
            ))
            
            # Add average line
            fig_yearly.add_trace(go.Scatter(
                x=daily_stats["datetime_display"],
                y=daily_stats["rh_avg"],
                mode="lines",
                name="Average Relative humidity",
                line=dict(color="#00a8ff", width=2),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.1f}%<extra></extra>",
            ))
            
            fig_yearly.update_layout(
                title="Annual Profile – Relative Humidity",
                xaxis_title="Day",
                yaxis_title="Relative Humidity (%)",
                hovermode="x unified",
                showlegend=True,
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=1.02,
                    xanchor="right",
                    x=1
                ),
                xaxis_rangeslider_visible=False,
                height=450,
                template="plotly_white",
                margin=dict(b=80),
            )
            
            st.plotly_chart(fig_yearly, use_container_width=True)
            
            # Calculate humidity metrics
            comfort_rh_lower = 40  # Comfort zone lower bound
            comfort_rh_upper = 60  # Comfort zone upper bound
            
            # Get min/max RH values
            rh_min = df["relative_humidity"].min()
            rh_max = df["relative_humidity"].max()
            rh_avg = df["relative_humidity"].mean()
            
            # Find date/time of peak RH
            max_rh_idx = df["relative_humidity"].idxmax()
            max_rh_row = df.loc[max_rh_idx]
            max_rh_date_str = max_rh_row["datetime"].strftime("%b %d")
            max_rh_hour = int(max_rh_row["hour"])
            
            # Comfort hours (40-60% RH)
            comfort_rh_hours = len(df[(df["relative_humidity"] >= comfort_rh_lower) & (df["relative_humidity"] <= comfort_rh_upper)])
            comfort_rh_percent = (comfort_rh_hours / len(df)) * 100
            
            # High humidity hours (> 60% RH)
            high_humidity_hrs = len(df[df["relative_humidity"] > 60])
            
            # Condensation risk hours (> 75% RH sustained)
            condensation_risk_hrs = len(df[df["relative_humidity"] > 75])
            
            # Low humidity hours (< 30% RH)
            low_humidity_hrs = len(df[df["relative_humidity"] < 30])
            
            # Mold risk hours (> 60% RH sustained for extended periods)
            mold_risk_hrs = len(df[df["relative_humidity"] > 60])
            
            # HVAC RH control (percentage within comfort band)
            hvac_rh_control = comfort_rh_percent
            
            # Overhumidification hours (> 70% RH)
            overhumidification_hrs = len(df[df["relative_humidity"] > 70])
            
            # Humidity metrics cards - Row 1
            rh_col1, rh_col2, rh_col3, rh_col4, rh_col5 = st.columns(5)
            
            with rh_col1:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #f59e0b;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Comfort 40-60%</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{comfort_rh_percent:.0f} %</div>
                        <div style="font-size: 11px; color: #718096;">Occupied RH Hrs</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col2:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #ef4444;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Peak RH (Occupied)</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_max:.1f} %</div>
                        <div style="font-size: 11px; color: #718096;">All year</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col3:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #8b5cf6;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #8b5cf6; text-transform: uppercase; letter-spacing: 0.5px;">High Humidity Hrs</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{high_humidity_hrs}</div>
                        <div style="font-size: 11px; color: #718096;">&gt; 60% RH</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col4:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #06b6d4;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">Condensation Risk Hrs</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{condensation_risk_hrs}</div>
                        <div style="font-size: 11px; color: #718096;">Surface Temp &lt; Dew Point</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col5:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #3b82f6;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Avg RH</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_avg:.1f} %</div>
                        <div style="font-size: 11px; color: #718096;"></div>
                    </div>
                """, unsafe_allow_html=True)
            
            # Row 2 - Additional metrics
            rh_col6, rh_col7, rh_col8, rh_col9, rh_col10 = st.columns(5)
            
            with rh_col6:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #f59e0b;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #f59e0b; text-transform: uppercase; letter-spacing: 0.5px;">Low Humidity Hrs</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{low_humidity_hrs}</div>
                        <div style="font-size: 11px; color: #718096;">&lt; 30% RH</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col7:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #ef4444;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #ef4444; text-transform: uppercase; letter-spacing: 0.5px;">Mold Risk Hrs</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{mold_risk_hrs}</div>
                        <div style="font-size: 11px; color: #718096;">&gt; 60% RH Sustained</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col8:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #06b6d4;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #06b6d4; text-transform: uppercase; letter-spacing: 0.5px;">HVAC RH Control</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{hvac_rh_control:.0f} %</div>
                        <div style="font-size: 11px; color: #718096;">Outside RH vs Inside RH</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col9:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #3b82f6;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #3b82f6; text-transform: uppercase; letter-spacing: 0.5px;">Overhumidification Hrs</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{overhumidification_hrs}</div>
                        <div style="font-size: 11px; color: #718096;">System Failure Indicator</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with rh_col10:
                st.markdown(f"""
                    <div style="
                        background: white;
                        padding: 16px;
                        border-radius: 8px;
                        border-left: 4px solid #0891b2;
                        box-shadow: 0 2px 4px rgba(0,0,0,0.08);
                        text-align: center;
                    ">
                        <div style="font-size: 11px; font-weight: 700; color: #0891b2; text-transform: uppercase; letter-spacing: 0.5px;">Min RH</div>
                        <div style="font-size: 26px; font-weight: 700; color: #2c3e50; margin: 8px 0;">{rh_min:.1f} %</div>
                        <div style="font-size: 11px; color: #718096;"></div>
                    </div>
                """, unsafe_allow_html=True)
        
        else:
            st.info("Sun Path analysis is not yet implemented.")

    # === TAB 2: MONTHLY TREND ===
    elif st.session_state.active_tab == "Monthly Trend":
        if selected_parameter == "Temperature":
            import plotly.graph_objects as go
            
            # Calculate monthly statistics
            monthly_stats = df.groupby("month").agg({
                "dry_bulb_temperature": ["min", "max", "mean"],
                "relative_humidity": ["min", "max", "mean"],
            }).reset_index()
            
            monthly_stats.columns = ["month", "temp_min", "temp_max", "temp_avg", "rh_min", "rh_max", "rh_avg"]
            month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
            monthly_stats["month_name"] = monthly_stats["month"].apply(lambda x: month_names[x-1])
            
            fig_monthly = go.Figure()
            
            # Get selected month range
            start_month = st.session_state.start_month_idx + 1
            end_month = st.session_state.end_month_idx + 1
            
            # GREYED OUT: Before selected range
            if start_month > 1:
                before_months = monthly_stats[monthly_stats["month"] < start_month]
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Period",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["temp_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected range
            active_months = monthly_stats[(monthly_stats["month"] >= start_month) & (monthly_stats["month"] <= end_month)]
            
            # Add temperature range (min/max) for active range
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["temp_max"],
                fill=None,
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["temp_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                name="Monthly Temperature Range",
                fillcolor="rgba(255, 173, 173, 0.4)",
                customdata=active_months["temp_max"],
                hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
            ))
            
            # Add average line for active range
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["temp_avg"],
                mode="lines+markers",
                name="Monthly Average Temperature",
                line=dict(color="#d32f2f", width=2),
                marker=dict(size=8),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
            ))
            
            # GREYED OUT: After selected range
            if end_month < 12:
                after_months = monthly_stats[monthly_stats["month"] > end_month]
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["temp_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            
            fig_monthly.update_layout(
                title="Monthly Temperature Trend",
                xaxis_title="Month",
                yaxis_title="Temperature (°C)",
                hovermode="x unified",
                showlegend=True,
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=1.02,
                    xanchor="right",
                    x=1
                ),
                height=450,
                template="plotly_white",
                margin=dict(b=80),
            )
            
            st.plotly_chart(fig_monthly, use_container_width=True)
            
            # Display monthly KPI metrics
            st.markdown("#### Monthly Temperature Summary")
            
            # Create a dataframe for monthly metrics
            kpi_data = monthly_stats[["month_name", "temp_min", "temp_max", "temp_avg"]].copy()
            kpi_data.columns = ["Month", "Min (°C)", "Max (°C)", "Avg (°C)"]
            
            # Display as a nice table
            st.dataframe(
                kpi_data,
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Min (°C)": st.column_config.NumberColumn(format="%.2f"),
                    "Max (°C)": st.column_config.NumberColumn(format="%.2f"),
                    "Avg (°C)": st.column_config.NumberColumn(format="%.2f"),
                }
            )
            
        elif selected_parameter == "Humidity":
            import plotly.graph_objects as go
            
            # Get selected month range for greying
            start_month = st.session_state.start_month_idx + 1
            end_month = st.session_state.end_month_idx + 1
            
            # Calculate monthly humidity statistics
            monthly_stats = df.groupby("month").agg({
                "relative_humidity": ["min", "max", "mean"],
            }).reset_index()
            
            monthly_stats.columns = ["month", "rh_min", "rh_max", "rh_avg"]
            month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
            monthly_stats["month_name"] = monthly_stats["month"].apply(lambda x: month_names[x-1])
            
            # Split data into before/active/after ranges
            before_months = monthly_stats[monthly_stats["month"] < start_month]
            active_months = monthly_stats[(monthly_stats["month"] >= start_month) & (monthly_stats["month"] <= end_month)]
            after_months = monthly_stats[monthly_stats["month"] > end_month]
            
            fig_monthly = go.Figure()
            
            # GREYED OUT: Before selected range
            if not before_months.empty:
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=[65] * len(before_months),
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=[30] * len(before_months),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Period",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=before_months["month_name"],
                    y=before_months["rh_avg"],
                    mode="lines",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected month range
            # Add humidity comfort band for active range
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=[65] * len(active_months),
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=[30] * len(active_months),
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="Humidity comfort band (30-65%)",
                fillcolor="rgba(128, 128, 128, 0.2)",
                hoverinfo="skip",
            ))
            
            # Add humidity range (min/max) for active range
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["rh_max"],
                fill=None,
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["rh_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                name="Monthly Humidity Range",
                fillcolor="rgba(0, 150, 255, 0.3)",
                customdata=active_months["rh_max"],
                hovertemplate="<b>%{x}</b><br>Min: %{y:.1f}%<br>Max: %{customdata:.1f}%<extra></extra>",
            ))
            
            # Add average line for active range
            fig_monthly.add_trace(go.Scatter(
                x=active_months["month_name"],
                y=active_months["rh_avg"],
                mode="lines+markers",
                name="Monthly Average Humidity",
                line=dict(color="#00a8ff", width=2),
                marker=dict(size=8),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.1f}%<extra></extra>",
            ))
            
            # GREYED OUT: After selected range
            if not after_months.empty:
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=[65] * len(after_months),
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=[30] * len(after_months),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_monthly.add_trace(go.Scatter(
                    x=after_months["month_name"],
                    y=after_months["rh_avg"],
                    mode="lines",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            fig_monthly.update_layout(
                title="Monthly Relative Humidity Trend",
                xaxis_title="Month",
                yaxis_title="Relative Humidity (%)",
                hovermode="x unified",
                showlegend=True,
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=1.02,
                    xanchor="right",
                    x=1
                ),
                height=450,
                template="plotly_white",
                margin=dict(b=80),
            )
            
            st.plotly_chart(fig_monthly, use_container_width=True)
            
            # Display monthly KPI metrics
            st.markdown("#### Monthly Humidity Summary")
            
            # Create a dataframe for monthly metrics
            kpi_data = monthly_stats[["month_name", "rh_min", "rh_max", "rh_avg"]].copy()
            kpi_data.columns = ["Month", "Min (%)", "Max (%)", "Avg (%)"]
            
            # Display as a nice table
            st.dataframe(
                kpi_data,
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Min (%)": st.column_config.NumberColumn(format="%.1f"),
                    "Max (%)": st.column_config.NumberColumn(format="%.1f"),
                    "Avg (%)": st.column_config.NumberColumn(format="%.1f"),
                }
            )
        
        else:
            st.info("Monthly trend analysis is not yet implemented for Sun Path.")

    # === TAB 3: DIURNAL PROFILE ===
    elif st.session_state.active_tab == "Diurnal Profile":
        import plotly.graph_objects as go
        
        if selected_parameter == "Temperature":
            # Create hourly averages for each month
            hourly_stats = df.groupby(["month", "hour"]).agg({
                "dry_bulb_temperature": ["min", "max", "mean"],
            }).reset_index()
            
            hourly_stats.columns = ["month", "hour", "temp_min", "temp_max", "temp_avg"]
            hourly_stats["month_name"] = hourly_stats["month"].apply(lambda x: ["Jan", "Feb", "Mar", "Apr", "May", "Jun", 
                                                                                  "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"][x-1])
            
            fig_diurnal = go.Figure()
            
            # Get hourly averages across ALL months for diurnal profile
            # Diurnal Profile uses TIME filter only, not month filter
            avg_hourly = hourly_stats.groupby("hour").agg({
                "temp_min": "min",
                "temp_max": "max",
                "temp_avg": "mean"
            }).reset_index()
            
            # Get selected hour range
            start_hour_sel, end_hour_sel = st.session_state.get("hour_range", (8, 18))
            
            # GREYED OUT: Before selected hours
            if start_hour_sel > 0:
                before_hours = avg_hourly[avg_hourly["hour"] < start_hour_sel]
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours["hour"],
                    y=before_hours["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours["hour"],
                    y=before_hours["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Hours",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours["hour"],
                    y=before_hours["temp_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected hours
            active_hours = avg_hourly[(avg_hourly["hour"] >= start_hour_sel) & (avg_hourly["hour"] <= end_hour_sel)]
            
            # Add temperature range for active hours
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours["hour"],
                y=active_hours["temp_max"],
                fill=None,
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours["hour"],
                y=active_hours["temp_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                name="Daily Range",
                fillcolor="rgba(255, 173, 173, 0.3)",
                customdata=active_hours["temp_max"],
                hovertemplate="<b>Hour %{x}:00</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
            ))
            
            # Add average for active hours
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours["hour"],
                y=active_hours["temp_avg"],
                mode="lines+markers",
                name="Average Temperature",
                line=dict(color="#d32f2f", width=2),
                marker=dict(size=6),
                hovertemplate="<b>Hour %{x}:00</b><br>Avg: %{y:.2f}°C<extra></extra>",
            ))
            
            # GREYED OUT: After selected hours
            if end_hour_sel < 23:
                after_hours = avg_hourly[avg_hourly["hour"] > end_hour_sel]
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours["hour"],
                    y=after_hours["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours["hour"],
                    y=after_hours["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours["hour"],
                    y=after_hours["temp_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            
            fig_diurnal.update_layout(
                title="Diurnal Temperature Profile",
                xaxis_title="Hour of Day",
                yaxis_title="Temperature (°C)",
                hovermode="x unified",
                showlegend=True,
                template="plotly_white",
                height=450,
            )
            
            st.plotly_chart(fig_diurnal, use_container_width=True)
        
        elif selected_parameter == "Humidity":
            # Create hourly humidity averages for each month
            hourly_humidity = df.groupby(["month", "hour"]).agg({
                "relative_humidity": ["min", "max", "mean"],
            }).reset_index()
            
            hourly_humidity.columns = ["month", "hour", "rh_min", "rh_max", "rh_avg"]
            
            fig_diurnal = go.Figure()
            
            # Get hourly averages across ALL months for diurnal profile
            avg_hourly_rh = hourly_humidity.groupby("hour").agg({
                "rh_min": "min",
                "rh_max": "max",
                "rh_avg": "mean"
            }).reset_index()
            
            # Get selected hour range
            start_hour_sel, end_hour_sel = st.session_state.get("hour_range", (8, 18))
            
            # GREYED OUT: Before selected hours
            if start_hour_sel > 0:
                before_hours_rh = avg_hourly_rh[avg_hourly_rh["hour"] < start_hour_sel]
                
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours_rh["hour"],
                    y=[65] * len(before_hours_rh),
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours_rh["hour"],
                    y=[30] * len(before_hours_rh),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Hours",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours_rh["hour"],
                    y=before_hours_rh["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours_rh["hour"],
                    y=before_hours_rh["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=before_hours_rh["hour"],
                    y=before_hours_rh["rh_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected hours
            active_hours_rh = avg_hourly_rh[(avg_hourly_rh["hour"] >= start_hour_sel) & (avg_hourly_rh["hour"] <= end_hour_sel)]
            
            # Add humidity comfort band for active hours
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours_rh["hour"],
                y=[65] * len(active_hours_rh),
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours_rh["hour"],
                y=[30] * len(active_hours_rh),
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="Comfort band (30-65%)",
                fillcolor="rgba(128, 128, 128, 0.2)",
                hoverinfo="skip",
            ))
            
            # Add humidity range for active hours
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours_rh["hour"],
                y=active_hours_rh["rh_max"],
                fill=None,
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours_rh["hour"],
                y=active_hours_rh["rh_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(0, 0, 255, 0)",
                name="Humidity Range",
                fillcolor="rgba(0, 150, 255, 0.3)",
                customdata=active_hours_rh["rh_max"],
                hovertemplate="<b>Hour %{x}:00</b><br>Min: %{y:.1f}%<br>Max: %{customdata:.1f}%<extra></extra>",
            ))
            
            # Add average for active hours
            fig_diurnal.add_trace(go.Scatter(
                x=active_hours_rh["hour"],
                y=active_hours_rh["rh_avg"],
                mode="lines+markers",
                name="Average Humidity",
                line=dict(color="#00a8ff", width=2),
                marker=dict(size=6),
                hovertemplate="<b>Hour %{x}:00</b><br>Avg: %{y:.1f}%<extra></extra>",
            ))
            
            # GREYED OUT: After selected hours
            if end_hour_sel < 23:
                after_hours_rh = avg_hourly_rh[avg_hourly_rh["hour"] > end_hour_sel]
                
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours_rh["hour"],
                    y=[65] * len(after_hours_rh),
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours_rh["hour"],
                    y=[30] * len(after_hours_rh),
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours_rh["hour"],
                    y=after_hours_rh["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours_rh["hour"],
                    y=after_hours_rh["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                
                fig_diurnal.add_trace(go.Scatter(
                    x=after_hours_rh["hour"],
                    y=after_hours_rh["rh_avg"],
                    mode="lines+markers",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    marker=dict(size=4),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            fig_diurnal.update_layout(
                title="Diurnal Humidity Profile",
                xaxis_title="Hour of Day",
                yaxis_title="Relative Humidity (%)",
                hovermode="x unified",
                showlegend=True,
                template="plotly_white",
                height=450,
            )
            
            st.plotly_chart(fig_diurnal, use_container_width=True)
        
        else:
            st.info("Sun Path analysis is not yet implemented.")
    
    # === TAB 4: COMFORT ANALYSIS ===
    elif st.session_state.active_tab == "Comfort Analysis":
        if selected_parameter == "Temperature":
            import plotly.graph_objects as go
            
            # Calculate day of year for selected month range for greying
            start_month_num = st.session_state.start_month_idx + 1
            end_month_num = st.session_state.end_month_idx + 1
            start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
            if end_month_num == 12:
                end_doy = 366
            else:
                end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
            
            # Create comfort analysis visualization
            fig_comfort = go.Figure()
            
            # GREYED OUT: Before selected range
            if start_doy > 1:
                before_data = daily_stats[daily_stats["doy"] < start_doy]
                fig_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Period",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                fig_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["temp_avg"],
                    mode="lines",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected range
            active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
            
            # Add comfort bands for active range
            fig_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_90_upper"],
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["comfort_90_lower"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="ASHRAE 90% acceptability",
                fillcolor="rgba(76, 175, 80, 0.4)",
                hoverinfo="skip",
            ))
            
            # Add temperature data for active range
            fig_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_max"],
                fill=None,
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(255, 0, 0, 0)",
                name="Daily Temperature Range",
                fillcolor="rgba(255, 173, 173, 0.3)",
                customdata=active_range["temp_max"],
                hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}°C<br>Max: %{customdata:.2f}°C<extra></extra>",
            ))
            
            # Add average for active range
            fig_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["temp_avg"],
                mode="lines",
                name="Average Temperature",
                line=dict(color="#d32f2f", width=2),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}°C<extra></extra>",
            ))
            
            # GREYED OUT: After selected range
            if end_doy < 365:
                after_data = daily_stats[daily_stats["doy"] > end_doy]
                if not after_data.empty:
                    fig_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["temp_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
            
            
            fig_comfort.update_layout(
                title="Comfort Analysis – ASHRAE Adaptive Comfort",
                xaxis_title="Day",
                yaxis_title="Temperature (°C)",
                hovermode="x unified",
                showlegend=True,
                template="plotly_white",
                height=450,
            )
            
            st.plotly_chart(fig_comfort, use_container_width=True)
        
        elif selected_parameter == "Humidity":
            import plotly.graph_objects as go
            
            # Calculate day of year for selected month range for greying
            start_month_num = st.session_state.start_month_idx + 1
            end_month_num = st.session_state.end_month_idx + 1
            start_doy = pd.to_datetime(f"2024-{start_month_num}-01").dayofyear
            if end_month_num == 12:
                end_doy = 366
            else:
                end_doy = pd.to_datetime(f"2024-{end_month_num+1}-01").dayofyear - 1
            
            # Create humidity comfort analysis visualization
            fig_humidity_comfort = go.Figure()
            
            # Define comfort zones for humidity
            comfort_upper = 60  # Upper comfort limit
            comfort_lower = 40  # Lower comfort limit
            
            # GREYED OUT: Before selected range
            if start_doy > 1:
                before_data = daily_stats[daily_stats["doy"] < start_doy]
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["rh_max"],
                    fill=None,
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    showlegend=False,
                    hoverinfo="skip",
                ))
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["rh_min"],
                    fill="tonexty",
                    mode="lines",
                    line_color="rgba(100, 100, 100, 0)",
                    fillcolor="rgba(180, 180, 180, 0.15)",
                    name="Unselected Period",
                    showlegend=True,
                    hoverinfo="skip",
                ))
                fig_humidity_comfort.add_trace(go.Scatter(
                    x=before_data["datetime_display"],
                    y=before_data["rh_avg"],
                    mode="lines",
                    line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                    showlegend=False,
                    hoverinfo="skip",
                ))
            
            # ACTIVE: Selected range
            active_range = daily_stats[(daily_stats["doy"] >= start_doy) & (daily_stats["doy"] <= end_doy)]
            
            # Add comfort band (40-60% RH) for active range
            fig_humidity_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=[comfort_upper] * len(active_range),
                fill=None,
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_humidity_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=[comfort_lower] * len(active_range),
                fill="tonexty",
                mode="lines",
                line_color="rgba(128, 128, 128, 0)",
                name="Comfort Band (40-60%)",
                fillcolor="rgba(76, 175, 80, 0.4)",
                hoverinfo="skip",
            ))
            
            # Add humidity data for active range
            fig_humidity_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["rh_max"],
                fill=None,
                mode="lines",
                line_color="rgba(0, 150, 255, 0)",
                showlegend=False,
                hoverinfo="skip",
            ))
            
            fig_humidity_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["rh_min"],
                fill="tonexty",
                mode="lines",
                line_color="rgba(0, 150, 255, 0)",
                name="Daily RH Range",
                fillcolor="rgba(0, 150, 255, 0.3)",
                customdata=active_range["rh_max"],
                hovertemplate="<b>%{x}</b><br>Min: %{y:.2f}%<br>Max: %{customdata:.2f}%<extra></extra>",
            ))
            
            # Add average for active range
            fig_humidity_comfort.add_trace(go.Scatter(
                x=active_range["datetime_display"],
                y=active_range["rh_avg"],
                mode="lines",
                name="Average RH",
                line=dict(color="#00a8ff", width=2),
                hovertemplate="<b>%{x}</b><br>Avg: %{y:.2f}%<extra></extra>",
            ))
            
            # GREYED OUT: After selected range
            if end_doy < 365:
                after_data = daily_stats[daily_stats["doy"] > end_doy]
                if not after_data.empty:
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["rh_max"],
                        fill=None,
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["rh_min"],
                        fill="tonexty",
                        mode="lines",
                        line_color="rgba(100, 100, 100, 0)",
                        fillcolor="rgba(180, 180, 180, 0.15)",
                        showlegend=False,
                        hoverinfo="skip",
                    ))
                    fig_humidity_comfort.add_trace(go.Scatter(
                        x=after_data["datetime_display"],
                        y=after_data["rh_avg"],
                        mode="lines",
                        line=dict(color="rgba(150, 150, 150, 0.4)", width=1, dash="dot"),
                        showlegend=False,
                        hoverinfo="skip",
                    ))
            
            fig_humidity_comfort.update_layout(
                title="Humidity Comfort Analysis – Optimal Range (40-60%)",
                xaxis_title="Day",
                yaxis_title="Relative Humidity (%)",
                hovermode="x unified",
                showlegend=True,
                template="plotly_white",
                height=450,
            )
            
            st.plotly_chart(fig_humidity_comfort, use_container_width=True)
        else:
            st.info("Comfort Analysis is available for Temperature and Humidity parameters.")
    
    # === TAB 5: ENERGY METRICS ===
    elif st.session_state.active_tab == "Energy Metrics":
        if selected_parameter == "Temperature":
            # Calculate energy metrics
            filtered_df = df[
                (df["datetime"].dt.date >= start_date) &
                (df["datetime"].dt.date <= end_date) &
                (df["hour"].between(start_hour, end_hour))
            ]
            
            if not filtered_df.empty:
                # HDD18 (Heating Degree Days at 18°C base)
                hdd18 = (18 - df["dry_bulb_temperature"]).clip(lower=0).sum()
                
                # CDD24 (Cooling Degree Days at 24°C base)
                cdd24 = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                
                # Additional energy metrics
                hdd18_filtered = (18 - filtered_df["dry_bulb_temperature"]).clip(lower=0).sum()
                cdd24_filtered = (df["dry_bulb_temperature"] - 24).clip(lower=0).sum()
                
                # Degree days by month
                monthly_hdd = df.groupby("month").apply(lambda x: (18 - x["dry_bulb_temperature"]).clip(lower=0).sum())
                monthly_cdd = df.groupby("month").apply(lambda x: (x["dry_bulb_temperature"] - 24).clip(lower=0).sum())
                
                month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                
                # Display metrics in cards
                st.markdown("#### Energy Performance Indicators")
                
                energy_col1, energy_col2, energy_col3, energy_col4 = st.columns(4)
                
                with energy_col1:
                    st.metric("HDD18 (Annual)", f"{hdd18:.0f}", "Heating Degree-Days")
                
                with energy_col2:
                    st.metric("CDD24 (Annual)", f"{cdd24:.0f}", "Cooling Degree-Days")
                
                with energy_col3:
                    st.metric("HDD18 (Period)", f"{hdd18_filtered:.0f}", "Heating Degree-Days")
                
                with energy_col4:
                    st.metric("CDD24 (Period)", f"{cdd24_filtered:.0f}", "Cooling Degree-Days")
                
                # Monthly breakdown chart
                import plotly.graph_objects as go
                from plotly.subplots import make_subplots
                
                fig_energy = make_subplots(specs=[[{"secondary_y": True}]])
                
                fig_energy.add_trace(
                    go.Bar(x=month_names, y=monthly_hdd.values, name="HDD18", marker_color="#2196F3"),
                    secondary_y=False,
                )
                
                fig_energy.add_trace(
                    go.Bar(x=month_names, y=monthly_cdd.values, name="CDD24", marker_color="#FF9800"),
                    secondary_y=False,
                )
                
                fig_energy.update_layout(
                    title="Monthly Degree-Days Distribution",
                    xaxis_title="Month",
                    yaxis_title="Degree-Days",
                    hovermode="x unified",
                    height=400,
                    barmode="group",
                )
                
                st.plotly_chart(fig_energy, use_container_width=True)
        
        elif selected_parameter == "Humidity":
            # Calculate humidity-related metrics
            filtered_df = df[
                (df["datetime"].dt.date >= start_date) &
                (df["datetime"].dt.date <= end_date) &
                (df["hour"].between(start_hour, end_hour))
            ]
            
            if not filtered_df.empty:
                # Calculate humidity metrics for full year
                high_humidity_annual = len(df[df["relative_humidity"] > 60])
                condensation_risk_annual = len(df[df["relative_humidity"] > 75])
                mold_risk_annual = len(df[df["relative_humidity"] > 60])
                low_humidity_annual = len(df[df["relative_humidity"] < 30])
                
                # Filtered period metrics
                high_humidity_filtered = len(filtered_df[filtered_df["relative_humidity"] > 60])
                condensation_risk_filtered = len(filtered_df[filtered_df["relative_humidity"] > 75])
                
                # Monthly breakdown
                monthly_high_rh = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 60]))
                monthly_condensation = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 75]))
                monthly_mold_risk = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] > 60]))
                monthly_low_rh = df.groupby("month").apply(lambda x: len(x[x["relative_humidity"] < 30]))
                
                month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
                
                # Display metrics in cards
                st.markdown("#### Humidity Performance Indicators")
                
                humidity_col1, humidity_col2, humidity_col3, humidity_col4 = st.columns(4)
                
                with humidity_col1:
                    st.metric("High RH Hrs (Annual)", f"{high_humidity_annual:.0f}", ">60% RH")
                
                with humidity_col2:
                    st.metric("Condensation Risk (Annual)", f"{condensation_risk_annual:.0f}", ">75% RH")
                
                with humidity_col3:
                    st.metric("High RH Hrs (Period)", f"{high_humidity_filtered:.0f}", ">60% RH")
                
                with humidity_col4:
                    st.metric("Condensation (Period)", f"{condensation_risk_filtered:.0f}", ">75% RH")
                
                # Monthly breakdown chart
                import plotly.graph_objects as go
                from plotly.subplots import make_subplots
                
                fig_humidity_energy = make_subplots(specs=[[{"secondary_y": True}]])
                
                fig_humidity_energy.add_trace(
                    go.Bar(x=month_names, y=monthly_high_rh.values, name="High RH (>60%)", marker_color="#0099ff"),
                    secondary_y=False,
                )
                
                fig_humidity_energy.add_trace(
                    go.Bar(x=month_names, y=monthly_condensation.values, name="Condensation Risk (>75%)", marker_color="#FF6B6B"),
                    secondary_y=False,
                )
                
                fig_humidity_energy.add_trace(
                    go.Scatter(x=month_names, y=monthly_low_rh.values, name="Low RH (<30%)", 
                              line=dict(color="#FFA500", width=2), mode="lines+markers"),
                    secondary_y=False,
                )
                
                fig_humidity_energy.update_layout(
                    title="Monthly Humidity Risk Distribution",
                    xaxis_title="Month",
                    yaxis_title="Hours",
                    hovermode="x unified",
                    height=400,
                    barmode="group",
                )
                
                st.plotly_chart(fig_humidity_energy, use_container_width=True)
        else:
            st.info("Energy Metrics is available for Temperature and Humidity parameters.")

# === GENERATE REPORT LOGIC ===
if st.session_state.get("generate_report", False):
    with st.spinner("Generating PowerPoint report..."):
        try:
            # Get filter parameters from left column
            start_month_num = st.session_state.start_month_idx + 1
            end_month_num = st.session_state.end_month_idx + 1
            year = df["datetime"].dt.year.iloc[0] if not df.empty else 2024
            
            start_date = pd.to_datetime(f"{year}-{start_month_num}-01").date()
            if end_month_num == 12:
                end_date = pd.to_datetime(f"{year}-12-31").date()
            else:
                end_date = (pd.to_datetime(f"{year}-{end_month_num+1}-01") - pd.Timedelta(days=1)).date()
            
            start_hour, end_hour = st.session_state.get("hour_range", (8, 18))
            
            # Generate the report
            report_bytes = generate_pptx_report(
                df, 
                start_date, 
                end_date, 
                start_hour, 
                end_hour, 
                selected_parameter
            )
            
            # Offer download
            st.download_button(
                label="📥 Download Report (PowerPoint)",
                data=report_bytes,
                file_name=f"Climate_Analysis_Report_{start_date.strftime('%Y%m%d')}_to_{end_date.strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_report"
            )
            
            # Reset the flag
            st.session_state.generate_report = False
            
        except Exception as e:
            st.error(f"❌ Failed to generate report: {str(e)}")
            st.session_state.generate_report = False

# Adding extra space at the bottom
st.markdown("<br><br>", unsafe_allow_html=True)